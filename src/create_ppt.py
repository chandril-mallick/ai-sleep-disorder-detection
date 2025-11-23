from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper to set title and content
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = title_text
        
        tf = content.text_frame
        for i, text in enumerate(content_text_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(24)
            p.space_after = Pt(14)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Samsung Sleep Health Predictor"
    subtitle.text = "AI-Driven Early Detection System\nSamsung Capstone Project"

    # 2. Problem Statement
    add_slide("Problem Statement: The Silent Epidemic", [
        "• Sleep disorders like Insomnia and Sleep Apnea affect millions globally.",
        "• Risks include cardiovascular disease, chronic fatigue, and reduced cognitive function.",
        "• Current diagnosis is expensive, time-consuming, and requires clinical sleep studies.",
        "• There is a critical gap for accessible, early screening tools."
    ])

    # 3. The Solution
    add_slide("The Solution: AI-Driven Dashboard", [
        "• An intelligent, web-based application for instant sleep health assessment.",
        "• Uses Machine Learning to analyze 11 key lifestyle and health markers.",
        "• Provides real-time risk classification: Healthy, Insomnia, or Sleep Apnea.",
        "• Accessible anywhere, empowering users to take proactive health steps."
    ])

    # 4. Technical Architecture
    add_slide("Technical Architecture", [
        "• Data Source: Sleep Health & Lifestyle Dataset (Blood Pressure, BMI, HRV, etc.)",
        "• Model: Random Forest Classifier (Optimized for tabular health data).",
        "• Backend: Python with Scikit-learn for inference.",
        "• Frontend: Streamlit for a responsive, interactive user interface.",
        "• Deployment: Lightweight and container-ready."
    ])

    # 5. Model Performance
    add_slide("Model Performance & Results", [
        "• Accuracy: 88% on unseen test data.",
        "• Precision: High precision in distinguishing Sleep Apnea from Healthy subjects.",
        "• Key Predictors: BMI Category, Systolic Blood Pressure, and Sleep Duration.",
        "• Validation: Rigorous train-test splitting and metric evaluation."
    ])

    # 6. Social Impact
    add_slide("Social Impact", [
        "• Democratizing Health: Brings clinical-grade screening to the user's home.",
        "• Prevention First: Early detection prevents long-term chronic conditions.",
        "• Public Safety: Reduces accidents caused by undiagnosed sleep deprivation.",
        "• Awareness: Educates users about the link between lifestyle (steps, stress) and sleep."
    ])

    # 7. Future Expansion
    add_slide("Future Expansion", [
        "• Wearable Integration: Sync directly with Samsung Galaxy Watch for real-time bio-data.",
        "• Deep Learning: Implement CNN-LSTM models for raw EEG/ECG signal analysis.",
        "• Telehealth Connection: One-click report sharing with sleep specialists.",
        "• Personalized Plans: AI-generated sleep hygiene recommendations based on risk."
    ])

    # Save
    output_path = "Samsung_Sleep_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()
